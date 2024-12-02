import streamlit as st
import os
import cohere
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

# Streamlit App Interface
st.title("Document-Based Q&A System")

# Define layout columns
col1, col2, col3 = st.columns(3)

# Left Column (Upload PDF)
with col1:
    st.subheader("Insert PDF")
    uploaded_pdf = st.file_uploader("Upload a PDF file", type="pdf")

# Center Column (Upload DOCX)
with col2:
    st.subheader("Insert Docx")
    uploaded_doc = st.file_uploader("Upload a DOC/DOCX file", type="docx")

# Right Column (Upload PPT)
with col3:
    st.subheader("Insert PPT")
    uploaded_ppt = st.file_uploader("Upload a PPT file", type="pptx")

# Ask a Question
st.subheader("Ask a Question")
question = st.text_input("Enter your question based on uploaded documents:")

if uploaded_pdf or uploaded_doc or uploaded_ppt:
    all_text = ""

    # Extract text from PDF
    if uploaded_pdf:
        pdf_text = ""
        pdf_reader = PdfReader(uploaded_pdf)
        for page in pdf_reader.pages:
            pdf_text += page.extract_text() or ""
        all_text += pdf_text + "\n"

    # Extract text from DOCX
    if uploaded_doc:
        doc_text = ""
        doc_file = Document(uploaded_doc)
        for paragraph in doc_file.paragraphs:
            doc_text += paragraph.text + "\n"
        all_text += doc_text + "\n"

    # Extract text from PPT
    if uploaded_ppt:
        ppt_text = ""
        ppt_file = Presentation(uploaded_ppt)
        for slide in ppt_file.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + "\n"
        all_text += ppt_text + "\n"

    # Process the text using LangChain and Cohere
    if question:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000, chunk_overlap=200, length_function=len, separators=['\n', '\n\n', ' ', '']
        )
        chunks = text_splitter.split_text(text=all_text)
        embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
        vectorstore = FAISS.from_texts(chunks, embedding=embeddings)
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 6})

        def format_docs(docs):
            return "\n\n".join(doc.page_content for doc in docs)

        retrieved_docs = retriever.get_relevant_documents(question)
        context = format_docs(retrieved_docs)

        prompt_template = """Answer the question as precisely as possible using the provided context. If the answer is
                            not contained in the context, say "answer not available in context." \n\n
                            Context: \n {context}\n
                            Question: \n {question} \n
                            Answer:"""
        prompt = PromptTemplate.from_template(template=prompt_template)
        full_prompt = prompt.format(context=context, question=question)

        os.environ['COHERE_API_KEY'] = "YOUR_COHERE_API_KEY"
        cohere_client = cohere.Client(os.getenv('COHERE_API_KEY'))
        response = cohere_client.generate(
            model='command',
            prompt=full_prompt,
            max_tokens=150,
            temperature=0.1
        )

        # Display the answer
        st.write("Answer:")
        st.write(response.generations[0].text.strip())
